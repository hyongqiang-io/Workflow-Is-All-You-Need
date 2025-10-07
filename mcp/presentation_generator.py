#!/usr/bin/env python3
"""
文生PPT MCP服务器 - Linus式简洁实现
Text-to-PPT MCP Server - Simple & Clean Implementation
"""

import json
import asyncio
import tempfile
import os
from datetime import datetime
from typing import Dict, Any, List, Optional
from pathlib import Path

from fastapi import FastAPI, HTTPException, Request
from pydantic import BaseModel
import uvicorn
import logging

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="Presentation Generator MCP Server",
    description="简洁的文生PPT MCP服务器",
    version="1.0.0"
)

# MCP工具定义 - 简单明了，做一件事并做好
PPT_TOOLS = [
    {
        "name": "create_presentation",
        "description": "根据文本内容创建PowerPoint演示文稿",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "演示文稿标题"
                },
                "slides": {
                    "type": "array",
                    "description": "幻灯片内容列表",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "幻灯片标题"},
                            "content": {"type": "string", "description": "幻灯片内容"},
                            "slide_type": {
                                "type": "string",
                                "enum": ["title", "content", "bullet_points"],
                                "default": "content",
                                "description": "幻灯片类型"
                            }
                        },
                        "required": ["title", "content"]
                    }
                },
                "template": {
                    "type": "string",
                    "enum": ["default", "modern", "minimal"],
                    "default": "default",
                    "description": "演示文稿模板"
                },
                "output_path": {
                    "type": "string",
                    "description": "输出文件路径（可选，默认为临时文件）"
                }
            },
            "required": ["title", "slides"]
        }
    },
    {
        "name": "parse_text_to_slides",
        "description": "将长文本智能解析为幻灯片结构",
        "parameters": {
            "type": "object",
            "properties": {
                "text": {
                    "type": "string",
                    "description": "需要解析的文本内容"
                },
                "max_slides": {
                    "type": "integer",
                    "default": 10,
                    "description": "最大幻灯片数量"
                }
            },
            "required": ["text"]
        }
    }
]

class MCPRequest(BaseModel):
    tool: str
    arguments: Dict[str, Any]

class PPTGenerator:
    """PPT生成器 - 核心逻辑，简单可靠"""

    def __init__(self):
        # 确保python-pptx可用
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.PP_ALIGN = PP_ALIGN
        except ImportError:
            raise ImportError("请安装python-pptx: pip install python-pptx")

    def create_presentation(self, title: str, slides: List[Dict],
                          template: str = "default",
                          output_path: Optional[str] = None) -> str:
        """创建PPT文件 - 核心功能"""

        # 创建演示文稿
        prs = self.Presentation()

        # 添加标题页
        title_slide_layout = prs.slide_layouts[0]  # 标题布局
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # 添加内容页
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            slide_content = slide_data.get("content", "")
            slide_type = slide_data.get("slide_type", "content")

            if slide_type == "bullet_points":
                # 项目符号布局
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_title

                # 处理项目符号内容
                content_placeholder = slide.placeholders[1]
                tf = content_placeholder.text_frame

                # 分割内容为列表项
                points = [line.strip() for line in slide_content.split('\n') if line.strip()]
                for i, point in enumerate(points):
                    if i == 0:
                        tf.text = f"• {point}"
                    else:
                        p = tf.add_paragraph()
                        p.text = f"• {point}"
                        p.level = 0
            else:
                # 内容布局
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_title
                slide.placeholders[1].text = slide_content

        # 确定输出路径
        if not output_path:
            temp_dir = tempfile.gettempdir()
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(temp_dir, f"presentation_{timestamp}.pptx")

        # 保存文件
        prs.save(output_path)
        logger.info(f"PPT已生成: {output_path}")

        return output_path

    def parse_text_to_slides(self, text: str, max_slides: int = 10) -> List[Dict]:
        """文本解析为幻灯片结构 - 简单而有效的解析"""

        # 按段落分割
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]

        slides = []
        current_slide = None

        for para in paragraphs[:max_slides * 2]:  # 防止过多内容
            # 检测标题（简单启发式规则）
            if (len(para) < 100 and
                (para.endswith(':') or para.endswith('：') or
                 para.isupper() or
                 any(keyword in para.lower() for keyword in ['第', '章', '部分', '概述', '总结']))):

                # 保存前一张幻灯片
                if current_slide:
                    slides.append(current_slide)
                    if len(slides) >= max_slides:
                        break

                # 开始新幻灯片
                current_slide = {
                    "title": para.replace(':', '').replace('：', ''),
                    "content": "",
                    "slide_type": "content"
                }
            else:
                # 内容段落
                if current_slide:
                    if current_slide["content"]:
                        current_slide["content"] += "\n\n"
                    current_slide["content"] += para
                else:
                    # 如果没有标题，创建一个默认标题
                    slides.append({
                        "title": "内容",
                        "content": para,
                        "slide_type": "content"
                    })

        # 添加最后一张幻灯片
        if current_slide and current_slide not in slides:
            slides.append(current_slide)

        # 如果没有解析出幻灯片，创建默认结构
        if not slides:
            slides = [{
                "title": "内容概述",
                "content": text[:500] + "..." if len(text) > 500 else text,
                "slide_type": "content"
            }]

        logger.info(f"解析出 {len(slides)} 张幻灯片")
        return slides

# 全局生成器实例
ppt_generator = PPTGenerator()

@app.get("/")
async def root():
    """根端点"""
    return {
        "message": "Presentation Generator MCP Server",
        "version": "1.0.0",
        "tools_count": len(PPT_TOOLS)
    }

@app.get("/health")
async def health_check():
    """健康检查"""
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "tools": [tool["name"] for tool in PPT_TOOLS]
    }

@app.get("/tools")
async def get_tools():
    """获取工具列表 - MCP标准接口"""
    return {"tools": PPT_TOOLS}

@app.post("/call")
async def call_tool(request: MCPRequest):
    """调用工具 - MCP标准接口"""
    tool_name = request.tool
    arguments = request.arguments

    logger.info(f"调用工具: {tool_name}")
    logger.debug(f"参数: {arguments}")

    try:
        if tool_name == "create_presentation":
            # 创建演示文稿
            title = arguments["title"]
            slides = arguments["slides"]
            template = arguments.get("template", "default")
            output_path = arguments.get("output_path")

            file_path = ppt_generator.create_presentation(
                title=title,
                slides=slides,
                template=template,
                output_path=output_path
            )

            return {
                "success": True,
                "result": {
                    "message": f"演示文稿创建成功",
                    "file_path": file_path,
                    "slides_count": len(slides) + 1,  # +1 for title slide
                    "template": template
                },
                "timestamp": datetime.now().isoformat()
            }

        elif tool_name == "parse_text_to_slides":
            # 解析文本为幻灯片
            text = arguments["text"]
            max_slides = arguments.get("max_slides", 10)

            slides = ppt_generator.parse_text_to_slides(text, max_slides)

            return {
                "success": True,
                "result": {
                    "message": f"文本解析完成，生成 {len(slides)} 张幻灯片",
                    "slides": slides,
                    "slides_count": len(slides)
                },
                "timestamp": datetime.now().isoformat()
            }

        else:
            raise ValueError(f"未知工具: {tool_name}")

    except Exception as e:
        logger.error(f"工具调用失败: {e}")
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }

# 兼容不同的调用格式
@app.post("/mcp/tools/call")
async def call_tool_mcp_format(request: Dict[str, Any]):
    """MCP标准格式调用"""
    # 转换为简单格式
    mcp_request = MCPRequest(
        tool=request.get("name", ""),
        arguments=request.get("arguments", {})
    )
    return await call_tool(mcp_request)

if __name__ == "__main__":
    print("🎯 启动文生PPT MCP服务器")
    print("📝 支持工具:")
    for tool in PPT_TOOLS:
        print(f"   - {tool['name']}: {tool['description']}")
    print(f"🌐 服务地址: http://0.0.0.0:8086")
    print(f"🔧 工具列表: curl http://localhost:8086/tools")
    print(f"📋 测试调用:")
    print('   curl -X POST http://localhost:8086/call \\')
    print('     -H "Content-Type: application/json" \\')
    print('     -d \'{"tool": "parse_text_to_slides", "arguments": {"text": "第一章：介绍\\n\\n这是一个测试文档。\\n\\n第二章：实现\\n\\n具体的实现步骤。"}}\'')

    uvicorn.run(
        app,
        host="0.0.0.0",
        port=8086,
        log_level="info"
    )